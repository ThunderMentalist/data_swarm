"""Attachment inventory + ingest helpers."""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any


def build_inventory(task_dir: Path, manifest_rows: list[dict]) -> dict[str, Any]:
    files = []
    for row in manifest_rows:
        files.append({
            "filename": row.get("filename", ""),
            "path": row.get("path", ""),
            "sha256": row.get("sha256", ""),
            "size_bytes": row.get("size_bytes", 0),
            "type": row.get("type", ""),
            "notes": row.get("notes", ""),
        })
    return {"task_dir": str(task_dir), "files": files}


def extract_text(task_dir: Path, row: dict, max_chars: int) -> tuple[str, str]:
    path = task_dir / row.get("path", "")
    ext = path.suffix.lower()
    if ext == ".pdf":
        return "", "PDF not ingested; provide txt/docx export."
    if ext in {".txt", ".md", ".csv", ".tsv", ".json", ".yaml", ".yml"}:
        try:
            text = path.read_text(encoding="utf-8", errors="ignore")
            return text[:max_chars], ""
        except OSError as exc:
            return "", str(exc)
    if ext == ".docx":
        try:
            import docx  # type: ignore
            doc = docx.Document(path)
            return "\n".join(p.text for p in doc.paragraphs)[:max_chars], ""
        except Exception:
            return "", "docx parser missing"
    if ext == ".xlsx":
        try:
            import openpyxl  # type: ignore
            wb = openpyxl.load_workbook(path, read_only=True)
            lines = []
            for ws in wb.worksheets:
                lines.append(f"# {ws.title}")
                for rowv in ws.iter_rows(min_row=1, max_row=5, values_only=True):
                    lines.append("\t".join("" if v is None else str(v) for v in rowv))
            return "\n".join(lines)[:max_chars], ""
        except Exception:
            return "", "xlsx parser missing"
    if ext == ".pptx":
        try:
            from pptx import Presentation  # type: ignore
            prs = Presentation(path)
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"# Slide {i}")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        lines.append(shape.text)
            return "\n".join(lines)[:max_chars], ""
        except Exception:
            return "", "pptx parser missing"
    return "", "unsupported file type"


def ingest_attachments(task_dir: Path, manifest_rows: list[dict], cfg: dict) -> tuple[dict, str, dict[str, str]]:
    max_chars = int(cfg.get("max_chars_per_file", 20000))
    max_files = int(cfg.get("max_files", 25))
    inventory = build_inventory(task_dir, manifest_rows[:max_files])
    extracted: dict[str, str] = {}
    lines = ["# Attachments Summary", ""]
    for row in manifest_rows[:max_files]:
        text, note = extract_text(task_dir, row, max_chars)
        fn = row.get("filename", "unknown")
        lines.append(f"- {fn}: {note or 'ingested'}")
        if text:
            extracted[fn] = text
    return inventory, "\n".join(lines) + "\n", extracted
